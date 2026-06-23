import os
import sys
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt

# Output directories
OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
SLIDES_DIR = os.path.join(OUTPUT_DIR, "generated_slides")
os.makedirs(SLIDES_DIR, exist_ok=True)

# PPTX output path
PPTX_PATH = os.path.join(os.path.dirname(OUTPUT_DIR), "Gridlock_AI_Presentation.pptx")

# Dimensions
WIDTH, HEIGHT = 1920, 1080

# Colors (Light Theme)
COLOR_BG = (248, 250, 252)        # Slate 50 (#F8FAFC)
COLOR_CARD_BG = (255, 255, 255)   # White (#FFFFFF)
COLOR_CARD_BORDER = (226, 232, 240) # Slate 200 (#E2E8F0)
COLOR_TEXT_PRIMARY = (15, 23, 42)  # Slate 900 (#0F172A)
COLOR_TEXT_SECONDARY = (71, 85, 105) # Slate 600 (#475569)
COLOR_TEXT_MUTED = (148, 163, 184) # Slate 400 (#94A3B8)

COLOR_BLUE = (37, 99, 235)        # Royal Blue (#2563EB)
COLOR_EMERALD = (16, 185, 129)    # Emerald 500 (#10B981)
COLOR_AMBER = (249, 115, 22)      # Orange 500 (#F97316)
COLOR_RED = (239, 68, 68)         # Red 500 (#EF4444)

# Load fonts
def get_fonts():
    font_paths = {
        "regular": "C:\\Windows\\Fonts\\segoeui.ttf",
        "bold": "C:\\Windows\\Fonts\\segoeuib.ttf",
        "semibold": "C:\\Windows\\Fonts\\seguisb.ttf"
    }
    
    # Fallback to default or Arial if Windows fonts don't exist
    fonts = {}
    for style, path in font_paths.items():
        if os.path.exists(path):
            fonts[style] = path
        else:
            # Fallback checks
            fallbacks = ["arial.ttf", "segoeui.ttf", "C:\\Windows\\Fonts\\arial.ttf"]
            found = False
            for fb in fallbacks:
                if os.path.exists(fb):
                    fonts[style] = fb
                    found = True
                    break
                # Try system paths on windows
                fb_path = os.path.join("C:\\Windows\\Fonts", fb)
                if os.path.exists(fb_path):
                    fonts[style] = fb_path
                    found = True
                    break
            if not found:
                fonts[style] = None
    return fonts

FONT_PATHS = get_fonts()

def load_font(style, size):
    path = FONT_PATHS.get(style)
    if path:
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            pass
    return ImageFont.load_default()

# Drawing helpers
def draw_card(draw, x, y, w, h, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER, border_width=1, radius=16):
    # Rounded rectangle card
    draw.rounded_rectangle(
        [(x, y), (x + w, y + h)],
        radius=radius,
        fill=bg_color,
        outline=border_color,
        width=border_width
    )

def draw_wrapped_text(draw, text, x, y, max_width, font, fill_color):
    words = text.split(' ')
    lines = []
    current_line = []
    
    for word in words:
        current_line.append(word)
        # Check size of line
        line_str = ' '.join(current_line)
        bbox = draw.textbbox((0, 0), line_str, font=font)
        line_w = bbox[2] - bbox[0]
        if line_w > max_width:
            current_line.pop()
            lines.append(' '.join(current_line))
            current_line = [word]
    if current_line:
        lines.append(' '.join(current_line))
        
    current_y = y
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=font)
        line_h = bbox[3] - bbox[1]
        draw.text((x, current_y), line, font=font, fill=fill_color)
        current_y += line_h + 8
    return current_y

def draw_centered_text(draw, text, cx, cy, font, fill_color):
    bbox = draw.textbbox((0, 0), text, font=font)
    w = bbox[2] - bbox[0]
    h = bbox[3] - bbox[1]
    draw.text((cx - w/2, cy - h/2 - 2), text, font=font, fill=fill_color)

def draw_header_footer(draw, slide_number, title=""):
    # Clear background
    draw.rectangle([(0, 0), (WIDTH, HEIGHT)], fill=COLOR_BG)
    
    if slide_number == 1 or slide_number == 11:
        # Title/Thank you slides have no header
        return
        
    # Draw top header
    # Flipkart logo text
    draw.text((80, 45), "Flipkart", font=load_font("bold", 30), fill=COLOR_BLUE)
    draw.text((205, 42), "|", font=load_font("regular", 32), fill=COLOR_TEXT_MUTED)
    draw.text((230, 46), "hackerearth", font=load_font("bold", 28), fill=(0, 0, 0))
    
    # Flipkart GRIDLOCK HACKATHON
    draw.text((1350, 46), "Flipkart", font=load_font("regular", 22), fill=COLOR_TEXT_PRIMARY)
    draw.text((1435, 46), "GRIDLOCK HACKATHON", font=load_font("bold", 22), fill=COLOR_TEXT_PRIMARY)
    
    # yellow badge "solve. traffic"
    badge_x, badge_y, badge_w, badge_h = 1690, 44, 150, 36
    draw.rounded_rectangle([(badge_x, badge_y), (badge_x + badge_w, badge_y + badge_h)], radius=8, fill=(254, 240, 138)) # light yellow
    draw_centered_text(draw, "solve. traffic", badge_x + badge_w/2, badge_y + badge_h/2, load_font("bold", 18), (113, 63, 18))
    
    # Divider line
    draw.line([(80, 105), (1840, 105)], fill=COLOR_CARD_BORDER, width=2)
    
    # Slide Title
    if title:
        draw.text((80, 135), title, font=load_font("bold", 38), fill=COLOR_TEXT_PRIMARY)
        
    # Footer
    draw.line([(80, 980), (1840, 980)], fill=COLOR_CARD_BORDER, width=1)
    draw.text((80, 1000), "Team Arohan  •  Flipkart GRIDLOCK Hackathon 2026", font=load_font("regular", 18), fill=COLOR_TEXT_SECONDARY)
    draw.text((900, 1000), "Gridlock.AI - Predictive Incident Dispatch Engine", font=load_font("semibold", 18), fill=COLOR_TEXT_MUTED)
    draw.text((1770, 1000), f"{slide_number:02d} / 11", font=load_font("bold", 18), fill=COLOR_TEXT_SECONDARY)
    
    # Bottom accent lines
    draw.rectangle([(0, 1072), (640, 1080)], fill=COLOR_BLUE)
    draw.rectangle([(640, 1072), (1280, 1080)], fill=COLOR_EMERALD)
    draw.rectangle([(1280, 1072), (1920, 1080)], fill=COLOR_AMBER)

# Icons Drawing Helpers
def draw_traffic_light_icon(draw, cx, cy):
    # Draw housing
    draw.rounded_rectangle([(cx - 20, cy - 40), (cx + 20, cy + 40)], radius=8, fill=(30, 41, 59))
    # Red light
    draw.ellipse([(cx - 10, cy - 30), (cx + 10, cy - 10)], fill=COLOR_RED)
    # Yellow light
    draw.ellipse([(cx - 10, cy - 10), (cx + 10, cy + 10)], fill=COLOR_AMBER)
    # Green light
    draw.ellipse([(cx - 10, cy + 10), (cx + 10, cy + 30)], fill=COLOR_EMERALD)

def draw_clock_icon(draw, cx, cy, color=COLOR_BLUE):
    # Outer circle
    draw.ellipse([(cx - 30, cy - 30), (cx + 30, cy + 30)], outline=color, width=4)
    # Clock hands
    draw.line([(cx, cy), (cx, cy - 18)], fill=color, width=4)
    draw.line([(cx, cy), (cx + 12, cy)], fill=color, width=4)

def draw_ambulance_icon(draw, cx, cy):
    # Body
    draw.rounded_rectangle([(cx - 30, cy - 15), (cx + 15, cy + 15)], radius=4, fill=COLOR_BG, outline=COLOR_RED, width=3)
    draw.rounded_rectangle([(cx + 15, cy - 5), (cx + 30, cy + 15)], radius=4, fill=COLOR_BG, outline=COLOR_RED, width=3)
    # Wheels
    draw.ellipse([(cx - 20, cy + 10), (cx - 10, cy + 20)], fill=(30, 41, 59))
    draw.ellipse([(cx + 15, cy + 10), (cx + 25, cy + 20)], fill=(30, 41, 59))
    # Red cross
    draw.line([(cx - 12, cy), (cx - 2, cy)], fill=COLOR_RED, width=3)
    draw.line([(cx - 7, cy - 5), (cx - 7, cy + 5)], fill=COLOR_RED, width=3)

def draw_alert_icon(draw, cx, cy, color=COLOR_RED):
    # Triangle alert
    points = [(cx, cy - 30), (cx - 30, cy + 20), (cx + 30, cy + 20)]
    draw.polygon(points, fill=(254, 226, 226), outline=color, width=4)
    # Exclamation
    draw.line([(cx, cy - 12), (cx, cy + 4)], fill=color, width=5)
    draw.ellipse([(cx - 3, cy + 10), (cx + 3, cy + 16)], fill=color)

def draw_brain_icon(draw, cx, cy):
    # Draw simple brain nodes
    draw.ellipse([(cx - 15, cy - 15), (cx + 15, cy + 15)], fill=(219, 234, 254), outline=COLOR_BLUE, width=3)
    draw.ellipse([(cx - 28, cy - 5), (cx - 18, cy + 5)], fill=(219, 234, 254), outline=COLOR_BLUE, width=3)
    draw.ellipse([(cx + 18, cy - 5), (cx + 28, cy + 5)], fill=(219, 234, 254), outline=COLOR_BLUE, width=3)
    draw.ellipse([(cx - 12, cy - 25), (cx - 2, cy - 15)], fill=(219, 234, 254), outline=COLOR_BLUE, width=3)
    draw.ellipse([(cx + 2, cy - 25), (cx + 12, cy - 15)], fill=(219, 234, 254), outline=COLOR_BLUE, width=3)

# ----------------- SLIDE GENERATOR FUNCTIONS -----------------

def make_slide_1():
    img = Image.new("RGB", (WIDTH, HEIGHT), COLOR_BG)
    draw = ImageDraw.Draw(img)
    
    # Flipkart Gridlock Banner Top Left
    draw.text((80, 45), "Flipkart", font=load_font("bold", 30), fill=COLOR_BLUE)
    draw.text((205, 42), "|", font=load_font("regular", 32), fill=COLOR_TEXT_MUTED)
    draw.text((230, 46), "hackerearth", font=load_font("bold", 28), fill=(0, 0, 0))
    
    # Top Right Banner
    draw.text((1350, 46), "Flipkart", font=load_font("regular", 22), fill=COLOR_TEXT_PRIMARY)
    draw.text((1435, 46), "GRIDLOCK HACKATHON", font=load_font("bold", 22), fill=COLOR_TEXT_PRIMARY)
    badge_x, badge_y, badge_w, badge_h = 1690, 44, 150, 36
    draw.rounded_rectangle([(badge_x, badge_y), (badge_x + badge_w, badge_y + badge_h)], radius=8, fill=(254, 240, 138))
    draw_centered_text(draw, "solve. traffic", badge_x + badge_w/2, badge_y + badge_h/2, load_font("bold", 18), (113, 63, 18))
    
    # Left Content Column (Hero Title)
    draw.text((80, 240), "SMART URBAN TRAFFIC", font=load_font("bold", 72), fill=COLOR_TEXT_PRIMARY)
    draw.text((80, 335), "MANAGEMENT & EMERGENCY", font=load_font("bold", 72), fill=COLOR_TEXT_PRIMARY)
    draw.text((80, 430), "RESPONSE SYSTEM", font=load_font("bold", 72), fill=COLOR_TEXT_PRIMARY)
    
    # Taglines
    draw.text((80, 550), "AI-POWERED  •  REAL-TIME  •  DATA-DRIVEN", font=load_font("semibold", 28), fill=COLOR_BLUE)
    
    # Problem Statement Box
    draw_card(draw, 80, 620, 780, 120, bg_color=(254, 242, 242), border_color=(252, 165, 165), border_width=2, radius=12)
    draw_alert_icon(draw, 130, 680, COLOR_RED)
    draw.text((190, 642), "Problem Statement Target:", font=load_font("semibold", 20), fill=COLOR_RED)
    draw.text((190, 675), "Event-Driven Congestion Mitigation (Planned & Unplanned Incidents)", font=load_font("bold", 22), fill=COLOR_TEXT_PRIMARY)
    
    # Team Info Card
    draw_card(draw, 80, 780, 780, 160, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER, border_width=1, radius=16)
    
    # Team name and details
    draw.text((120, 805), "TEAM AROHAN", font=load_font("bold", 24), fill=COLOR_BLUE)
    
    draw.text((120, 855), "Team Leader:", font=load_font("semibold", 18), fill=COLOR_TEXT_SECONDARY)
    draw.text((120, 885), "Nikita Pawar", font=load_font("bold", 20), fill=COLOR_TEXT_PRIMARY)
    
    draw.text((400, 855), "Teammates:", font=load_font("semibold", 18), fill=COLOR_TEXT_SECONDARY)
    draw.text((400, 885), "Jay Bhandarkar  •  Prajwal Kate  •  Riya Lad", font=load_font("bold", 20), fill=COLOR_TEXT_PRIMARY)
    
    # Right Side Graphic (Command Center Mockup HUD + Grid Nodes)
    draw_card(draw, 930, 240, 910, 700, bg_color=(15, 23, 42), border_color=(51, 65, 85), border_width=2, radius=24) # Dark dashboard card
    
    # Dashboard Header
    draw.text((970, 275), "GRIDLOCK.AI - COMMAND CENTER", font=load_font("bold", 22), fill=(248, 250, 252))
    # Pulse indicator
    draw.ellipse([(1800, 280), (1816, 296)], fill=COLOR_RED)
    draw.text((1720, 278), "LIVE", font=load_font("semibold", 16), fill=COLOR_RED)
    
    # Metrics row inside dark card
    metrics = [
        ("ACTIVE INCIDENTS", "128", "+12%", COLOR_RED),
        ("CRITICAL severity", "12", "High Risk", COLOR_AMBER),
        ("DISPATCHED OFFICERS", "45", "On-Field", COLOR_EMERALD)
    ]
    for i, (label, val, trend, color) in enumerate(metrics):
        mx = 970 + i * 270
        draw_card(draw, mx, 330, 250, 110, bg_color=(30, 41, 59), border_color=(71, 85, 105), border_width=1, radius=12)
        draw.text((mx + 20, 350), label, font=load_font("semibold", 12), fill=COLOR_TEXT_MUTED)
        draw.text((mx + 20, 375), val, font=load_font("bold", 32), fill=(248, 250, 252))
        draw.text((mx + 200, 385), trend, font=load_font("regular", 11), fill=color)
        
    # Map area grid representation
    map_x, map_y, map_w, map_h = 970, 470, 830, 430
    draw_card(draw, map_x, map_y, map_w, map_h, bg_color=(13, 18, 30), border_color=(30, 41, 59), border_width=1, radius=16)
    
    # Draw simulated road network grid
    import random
    random.seed(42)
    nodes = []
    for _ in range(12):
        nodes.append((random.randint(map_x + 60, map_x + map_w - 60), random.randint(map_y + 60, map_y + map_h - 60)))
    
    # Draw links
    for i in range(len(nodes)):
        for j in range(i+1, len(nodes)):
            dist = ((nodes[i][0] - nodes[j][0])**2 + (nodes[i][1] - nodes[j][1])**2)**0.5
            if dist < 220:
                draw.line([nodes[i], nodes[j]], fill=(30, 41, 59), width=2)
                
    # Draw hotspot circles and markers
    for idx, node in enumerate(nodes):
        if idx in [1, 5, 8]:
            # Critical hotspot
            draw.ellipse([(node[0] - 30, node[1] - 30), (node[0] + 30, node[1] + 30)], fill=(239, 68, 68, 80), outline=COLOR_RED, width=2)
            draw.ellipse([(node[0] - 6, node[1] - 6), (node[0] + 6, node[1] + 6)], fill=(255, 255, 255))
        elif idx in [3, 7]:
            # High severity
            draw.ellipse([(node[0] - 22, node[1] - 22), (node[0] + 22, node[1] + 22)], fill=(249, 115, 22, 80), outline=COLOR_AMBER, width=2)
            draw.ellipse([(node[0] - 5, node[1] - 5), (node[0] + 5, node[1] + 5)], fill=(255, 255, 255))
        else:
            # Normal flows
            draw.ellipse([(node[0] - 12, node[1] - 12), (node[0] + 12, node[1] + 12)], fill=(16, 185, 129, 60), outline=COLOR_EMERALD, width=1)
            draw.ellipse([(node[0] - 4, node[1] - 4), (node[0] + 4, node[1] + 4)], fill=(255, 255, 255))

    # Add a Map overlay label
    draw_card(draw, map_x + 20, map_y + 20, 240, 45, bg_color=(30, 41, 59), border_color=(71, 85, 105), border_width=1, radius=6)
    draw_centered_text(draw, "BENGALURU HEATMAP HUD", map_x + 140, map_y + 42, load_font("bold", 12), (248, 250, 252))

    # Bottom accent line
    draw.rectangle([(0, 1072), (640, 1080)], fill=COLOR_BLUE)
    draw.rectangle([(640, 1072), (1280, 1080)], fill=COLOR_EMERALD)
    draw.rectangle([(1280, 1072), (1920, 1080)], fill=COLOR_AMBER)

    img.save(os.path.join(SLIDES_DIR, "slide_01.png"))

def make_slide_2():
    img = Image.new("RGB", (WIDTH, HEIGHT), COLOR_BG)
    draw = ImageDraw.Draw(img)
    draw_header_footer(draw, 2, "Problem Statement & Existing Challenges")
    
    # Left Card: Problem Statement
    draw_card(draw, 80, 210, 840, 730, radius=20)
    
    # Icon and Title
    draw_alert_icon(draw, 160, 290, COLOR_BLUE)
    draw.text((220, 265), "The Problem Statement", font=load_font("bold", 28), fill=COLOR_BLUE)
    
    p_text_1 = "Rapid, unchecked urbanization has made traffic gridlocks a permanent fixture in modern metropolitan areas like Bengaluru."
    p_text_2 = "While long-term infrastructure improvements are capital-intensive and slow, emergency vehicles (like ambulances) and daily commuters face drastic delay penalties right now."
    p_text_3 = "Gridlocks are heavily compounded by event-driven incidents: both planned occurrences (metro construction, VIP road closures) and unplanned emergencies (heavy vehicle breakdowns, localized flooding, accidents)."
    p_text_4 = "Current city systems operate reactively—relying on fixed static signal cycles and dispatching traffic personnel only after long queues form, leading to massive productivity losses and excessive environmental emissions."

    y = 350
    for p in [p_text_1, p_text_2, p_text_3, p_text_4]:
        y = draw_wrapped_text(draw, p, 120, y, 760, load_font("regular", 20), COLOR_TEXT_PRIMARY) + 15
        
    # Traffic metrics at bottom
    draw_card(draw, 120, 770, 760, 130, bg_color=(241, 245, 249), border_color=COLOR_CARD_BORDER, radius=12)
    stats = [
        ("30-40%", "Increase in Peak Travel Times"),
        ("15-20 Min", "Emergency Dispatch Delays"),
        ("Millions", "Liters Wasted Annually")
    ]
    for i, (num, desc) in enumerate(stats):
        cx = 120 + i * 253 + 126
        draw_centered_text(draw, num, cx, 810, load_font("bold", 30), COLOR_BLUE)
        draw_centered_text(draw, desc, cx, 855, load_font("semibold", 14), COLOR_TEXT_SECONDARY)

    # Right Card: Existing Challenges
    draw_card(draw, 960, 210, 880, 730, radius=20)
    draw.text((1000, 250), "Key Structural Vulnerabilities", font=load_font("bold", 26), fill=COLOR_TEXT_PRIMARY)
    
    challenges = [
        ("Static Signal Control", "Fixed signal durations fail to dynamically adapt to sudden surges caused by localized incidents.", draw_traffic_light_icon),
        ("Delayed Emergency Routing", "Ambulances and emergency services get trapped in bottlenecks with zero priority preemption.", draw_ambulance_icon),
        ("Reactive Action Plans", "Dispatch control centers respond to gridlocks after they form, instead of predicting the impact.", draw_clock_icon),
        ("Disconnected Data Silos", "Junction video cameras, GPS fleets, weather reports, and operational logs operate independently.", draw_alert_icon),
        ("Economic & Environmental Drag", "Idle engine runtimes at major intersections result in massive fuel waste and toxic carbon emissions.", draw_alert_icon)
    ]
    
    cy = 310
    for i, (title, desc, icon_func) in enumerate(challenges):
        draw_card(draw, 1000, cy, 800, 105, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER, radius=12)
        if icon_func == draw_alert_icon:
            icon_func(draw, 1050, cy + 52, COLOR_AMBER)
        else:
            icon_func(draw, 1050, cy + 52)
        draw.text((1110, cy + 18), title, font=load_font("bold", 18), fill=COLOR_TEXT_PRIMARY)
        draw_wrapped_text(draw, desc, 1110, cy + 45, 660, load_font("regular", 14), COLOR_TEXT_SECONDARY)
        cy += 120

    img.save(os.path.join(SLIDES_DIR, "slide_02.png"))

def make_slide_3():
    img = Image.new("RGB", (WIDTH, HEIGHT), COLOR_BG)
    draw = ImageDraw.Draw(img)
    draw_header_footer(draw, 3, "Opportunity & Proposed Solution")
    
    # Left Card: The Opportunity
    draw_card(draw, 80, 210, 840, 730, radius=20)
    
    draw_alert_icon(draw, 160, 290, COLOR_AMBER)
    draw.text((220, 265), "The Opportunity", font=load_font("bold", 28), fill=COLOR_AMBER)
    
    opp_points = [
        ("Identify Incident Patterns Proactively", "Predict the severity and ripple effect of traffic disruptions immediately when an incident is logged by operators or sensors, moving the city from a reactive posture to predictive orchestration."),
        ("Automate Dispatch Protocols", "Bridge the gap between incident discovery and field response by instantly calculating required field resources (officers, physical barricades, diversions) using a standardized AI-guided engine."),
        ("Establish Dynamic Hotspot Scoring", "Aggregate multiple historical events across city corridors to identify chronic congestion nodes, enabling structural planning based on continuous empirical incident footprints."),
        ("Maximize System Availability", "Ensure operational command centers stay up 24/7/365 through cross-platform database replication and robust local fallback systems that operate even during network isolation.")
    ]
    
    y = 340
    for title, desc in opp_points:
        draw.text((120, y), "•  " + title, font=load_font("bold", 18), fill=COLOR_TEXT_PRIMARY)
        y = draw_wrapped_text(draw, desc, 140, y + 25, 730, load_font("regular", 15), COLOR_TEXT_SECONDARY) + 20

    # Right Card: The Solution
    draw_card(draw, 960, 210, 880, 730, radius=20)
    
    draw_brain_icon(draw, 1040, 290)
    draw.text((1100, 265), "Our Solution: Gridlock.AI", font=load_font("bold", 28), fill=COLOR_BLUE)
    
    sol_points = [
        ("AI Event Triage", "Utilizes an in-process CatBoost model to evaluate high-cardinality geographic features and cyclical dates, calculating an incident severity rating in under 10 milliseconds.", COLOR_BLUE),
        ("Dynamic Severity Index", "Translates discrete classes into a continuous Severity Score (0.0 to 1.0) using class probability weights, allowing granular prioritization between multiple critical incidents.", COLOR_BLUE),
        ("Resource Optimization Engine", "A rule-based recommendation matrix that auto-allocates required field equipment (officers, barricades) based on the calculated severity band, speeding up response times.", COLOR_BLUE),
        ("Interactive HUD Visualizations", "Integrates Leaflet map HUDs with custom filters and Recharts analytics dashboards to provide control room operators with real-time situational awareness.", COLOR_BLUE),
        ("Dual-Tier Storage Sync", "Synchronizes telemetry with cloud PostgreSQL (Supabase) while maintaining a local SQLite database engine that activates seamlessly to keep systems functional offline.", COLOR_BLUE)
    ]
    
    cy = 345
    for title, desc, col in sol_points:
        draw_card(draw, 1000, cy, 800, 100, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER, radius=12)
        draw.text((1025, cy + 18), title, font=load_font("bold", 18), fill=COLOR_BLUE)
        draw_wrapped_text(draw, desc, 1025, cy + 45, 750, load_font("regular", 14), COLOR_TEXT_SECONDARY)
        cy += 115

    img.save(os.path.join(SLIDES_DIR, "slide_03.png"))

def make_slide_4():
    img = Image.new("RGB", (WIDTH, HEIGHT), COLOR_BG)
    draw = ImageDraw.Draw(img)
    draw_header_footer(draw, 4, "Unique Selling Proposition (USP) & Key Features")
    
    # Left Card: Why Gridlock.AI (USP)
    draw_card(draw, 80, 210, 840, 730, radius=20)
    draw.text((120, 260), "Why Gridlock.AI? (Our USP)", font=load_font("bold", 28), fill=COLOR_BLUE)
    
    usps = [
        ("Hybrid AI & Rules Architecture", "We combine the predictive capability of Machine Learning (to classify incident severity) with a deterministic rules engine (to output physical resource recommendations), creating a system that is both intelligent and operationally safe."),
        ("High-Cardinality Spatial Native", "Unlike naive models that suffer from sparse representations, our CatBoost engine uses Ordered Target Encoding to process high-cardinality geographical attributes (junctions, zones, corridors) natively, maintaining high prediction accuracy."),
        ("Zero-Downtime Local Fallback", "Our double-tiered DB wrapper automatically senses remote network state. If the cloud database (Supabase) is unavailable, it redirects all writes and queries to local SQLite, executing a silent sync when connection is restored."),
        ("Continuous Impact Metrics", "By calculating a continuous 0.0 to 1.0 index from raw ML probability spreads, we let controllers prioritize between two high-priority events, resolving the limitations of simple classification.")
    ]
    
    y = 330
    for title, desc in usps:
        draw.text((120, y), "★  " + title, font=load_font("bold", 20), fill=COLOR_TEXT_PRIMARY)
        y = draw_wrapped_text(draw, desc, 145, y + 30, 730, load_font("regular", 15), COLOR_TEXT_SECONDARY) + 20

    # Right Card: Key Features Grid
    draw_card(draw, 960, 210, 880, 730, radius=20)
    draw.text((1000, 260), "Key Operational Features", font=load_font("bold", 28), fill=COLOR_TEXT_PRIMARY)
    
    features = [
        ("AI Event Triage", "Predicts priority (Low, Med, High, Critical) and road closure needs, generating confidence metrics for verification.", COLOR_BLUE, draw_brain_icon),
        ("Operations Dashboard", "Displays live Map HUDs showing active incident hotspots with proportional visual warning markers.", COLOR_EMERALD, draw_clock_icon),
        ("Advanced Analytics", "Provides charts on incident trends, clearance intervals, and cause distributions powered by Recharts.", COLOR_AMBER, draw_alert_icon),
        ("DB Sync & Fallback", "Keeps command center resilient with automated SQLite failovers and PostgreSQL sync routines.", COLOR_RED, draw_ambulance_icon)
    ]
    
    grid_coords = [
        (1000, 330),
        (1420, 330),
        (1000, 610),
        (1420, 610)
    ]
    
    for idx, (title, desc, col, icon_func) in enumerate(features):
        gx, gy = grid_coords[idx]
        draw_card(draw, gx, gy, 400, 250, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER, radius=16)
        icon_func(draw, gx + 60, gy + 60)
        draw.text((gx + 120, gy + 45), title, font=load_font("bold", 20), fill=col)
        draw_wrapped_text(draw, desc, gx + 30, gy + 110, 340, load_font("regular", 14), COLOR_TEXT_SECONDARY)

    img.save(os.path.join(SLIDES_DIR, "slide_04.png"))

def make_slide_5():
    img = Image.new("RGB", (WIDTH, HEIGHT), COLOR_BG)
    draw = ImageDraw.Draw(img)
    draw_header_footer(draw, 5, "User Journey & Process Flow Diagram")
    
    draw_card(draw, 80, 210, 1760, 730, radius=24)
    draw.text((120, 260), "Incident Life-Cycle & Operations Flow", font=load_font("bold", 28), fill=COLOR_TEXT_PRIMARY)
    
    steps = [
        ("1. Simulation Input", "Operator enters incident parameters on Next.js form: coordinates, corridor, zone, vehicle types, start time."),
        ("2. API Ingestion", "FastAPI backend handles payload, validates schema via Pydantic, and extracts temporal cyclical variables."),
        ("3. CatBoost Inference", "Model loads `remediated_traffic_model.cbm`, processes categorical strings, and predicts class probabilities."),
        ("4. Impact Score", "API computes continuous Severity Index (0.0 to 1.0) and determines the required officers, barricades, and diversions."),
        ("5. Dynamic Hotspot Sync", "Database client attempts PostgreSQL (Supabase) insert; triggers SQLite fallback if offline. Performs junction score UPSERT."),
        ("6. Map HUD Refresh", "Next.js pulls updated hotspot scores and refreshes Leaflet Map HUD with proportional color-coded warning circles.")
    ]
    
    coords = [
        (150, 360, 480, 180),
        (720, 360, 480, 180),
        (1290, 360, 480, 180),
        (1290, 640, 480, 180),
        (720, 640, 480, 180),
        (150, 640, 480, 180)
    ]
    
    colors = [COLOR_BLUE, COLOR_BLUE, COLOR_AMBER, COLOR_AMBER, COLOR_EMERALD, COLOR_EMERALD]
    
    for i, (title, desc) in enumerate(steps):
        cx, cy, cw, ch = coords[i]
        draw_card(draw, cx, cy, cw, ch, bg_color=COLOR_CARD_BG, border_color=colors[i], border_width=2, radius=12)
        
        draw.ellipse([(cx + 20, cy + 20), (cx + 55, cy + 55)], fill=colors[i])
        draw_centered_text(draw, str(i+1), cx + 37.5, cy + 37.5, load_font("bold", 18), (255, 255, 255))
        
        draw.text((cx + 70, cy + 25), title, font=load_font("bold", 18), fill=COLOR_TEXT_PRIMARY)
        draw_wrapped_text(draw, desc, cx + 25, cy + 70, cw - 50, load_font("regular", 13), COLOR_TEXT_SECONDARY)
        
    # Draw arrows
    draw.line([(630, 450), (720, 450)], fill=COLOR_TEXT_MUTED, width=4)
    draw.polygon([(720, 450), (705, 440), (705, 460)], fill=COLOR_TEXT_MUTED)
    
    draw.line([(1200, 450), (1290, 450)], fill=COLOR_TEXT_MUTED, width=4)
    draw.polygon([(1290, 450), (1275, 440), (1275, 460)], fill=COLOR_TEXT_MUTED)
    
    draw.line([(1530, 540), (1530, 640)], fill=COLOR_TEXT_MUTED, width=4)
    draw.polygon([(1530, 640), (1520, 625), (1540, 625)], fill=COLOR_TEXT_MUTED)
    
    draw.line([(1290, 730), (1200, 730)], fill=COLOR_TEXT_MUTED, width=4)
    draw.polygon([(1200, 730), (1215, 720), (1215, 740)], fill=COLOR_TEXT_MUTED)
    
    draw.line([(720, 730), (630, 730)], fill=COLOR_TEXT_MUTED, width=4)
    draw.polygon([(630, 730), (645, 720), (645, 740)], fill=COLOR_TEXT_MUTED)

    img.save(os.path.join(SLIDES_DIR, "slide_05.png"))

def make_slide_6():
    img = Image.new("RGB", (WIDTH, HEIGHT), COLOR_BG)
    draw = ImageDraw.Draw(img)
    draw_header_footer(draw, 6, "Wireframes: Command Center Dashboard")
    
    # Left Card: Map HUD
    draw_card(draw, 80, 210, 840, 730, radius=20)
    draw.text((120, 250), "Real-Time Traffic Operations Map", font=load_font("bold", 24), fill=COLOR_TEXT_PRIMARY)
    
    map_x, map_y, map_w, map_h = 120, 300, 760, 430
    draw_card(draw, map_x, map_y, map_w, map_h, bg_color=(15, 23, 42), border_color=(51, 65, 85), border_width=1, radius=12)
    
    for gx in range(map_x + 50, map_x + map_w, 100):
        draw.line([(gx, map_y), (gx, map_y + map_h)], fill=(30, 41, 59), width=1)
    for gy in range(map_y + 50, map_y + map_h, 100):
        draw.line([(map_x, gy), (map_x + map_w, gy)], fill=(30, 41, 59), width=1)
        
    node_a = (map_x + 200, map_y + 120)
    node_b = (map_x + 400, map_y + 240)
    node_c = (map_x + 600, map_y + 350)
    
    draw.line([node_a, node_b], fill=(51, 65, 85), width=4)
    draw.line([node_b, node_c], fill=(51, 65, 85), width=4)
    
    draw.ellipse([(node_a[0] - 15, node_a[1] - 15), (node_a[0] + 15, node_a[1] + 15)], fill=(16, 185, 129, 100), outline=COLOR_EMERALD, width=2)
    draw.ellipse([(node_b[0] - 25, node_b[1] - 25), (node_b[0] + 25, node_b[1] + 25)], fill=(249, 115, 22, 100), outline=COLOR_AMBER, width=2)
    draw.ellipse([(node_c[0] - 40, node_c[1] - 40), (node_c[0] + 40, node_c[1] + 40)], fill=(239, 68, 68, 100), outline=COLOR_RED, width=3)
    
    draw.text((node_a[0] + 20, node_a[1] - 10), "Hebbal (Low)", font=load_font("semibold", 12), fill=(248, 250, 252))
    draw.text((node_b[0] + 30, node_b[1] - 10), "Central (Medium)", font=load_font("semibold", 12), fill=(248, 250, 252))
    draw.text((node_c[0] - 120, node_c[1] - 50), "Silk Board (Critical Hotspot)", font=load_font("bold", 13), fill=(248, 250, 252))
    
    draw_card(draw, 120, 750, 760, 160, bg_color=COLOR_BG, border_color=COLOR_CARD_BORDER, radius=12)
    draw.text((150, 770), "Interactive Operations Features:", font=load_font("bold", 16), fill=COLOR_TEXT_PRIMARY)
    draw.text((150, 800), "• Click hotspots to view historical average incident duration and scores.", font=load_font("regular", 14), fill=COLOR_TEXT_SECONDARY)
    draw.text((150, 830), "• Color scales dynamically: Green (0.0-0.3), Orange (0.3-0.6), Red (0.6-1.0).", font=load_font("regular", 14), fill=COLOR_TEXT_SECONDARY)
    draw.text((150, 860), "• Filter map overlays by zone (South, North, East) or incident cause.", font=load_font("regular", 14), fill=COLOR_TEXT_SECONDARY)

    # Right Card: Event Explorer
    draw_card(draw, 960, 210, 880, 730, radius=20)
    draw.text((1000, 250), "Incident Database & Live Feed", font=load_font("bold", 24), fill=COLOR_TEXT_PRIMARY)
    
    table_x, table_y, table_w, table_h = 1000, 300, 800, 430
    draw_card(draw, table_x, table_y, table_w, table_h, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER, radius=12)
    
    draw.rectangle([(table_x, table_y), (table_x + table_w, table_y + 45)], fill=(241, 245, 249))
    draw.text((table_x + 20, table_y + 12), "INCIDENT ID", font=load_font("bold", 13), fill=COLOR_TEXT_SECONDARY)
    draw.text((table_x + 180, table_y + 12), "JUNCTION", font=load_font("bold", 13), fill=COLOR_TEXT_SECONDARY)
    draw.text((table_x + 360, table_y + 12), "CAUSE", font=load_font("bold", 13), fill=COLOR_TEXT_SECONDARY)
    draw.text((table_x + 500, table_y + 12), "SEVERITY", font=load_font("bold", 13), fill=COLOR_TEXT_SECONDARY)
    draw.text((table_x + 650, table_y + 12), "RESOURCES", font=load_font("bold", 13), fill=COLOR_TEXT_SECONDARY)
    
    events_mock = [
        ("FK-88412", "Silk Board Junc", "Accident", "CRITICAL", "15 Off, 30 Bar", COLOR_RED),
        ("FK-88405", "Hebbal Flyover", "Water Logging", "HIGH", "10 Off, 20 Bar", COLOR_AMBER),
        ("FK-88392", "Tin Factory Junc", "Vehicle Breakdown", "MEDIUM", "5 Off, 10 Bar", COLOR_BLUE),
        ("FK-88381", "Bannerghatta Rd", "Road Construction", "MEDIUM", "5 Off, 10 Bar", COLOR_BLUE),
        ("FK-88370", "Richmond Circle", "VIP Movement", "LOW", "2 Off, 0 Bar", COLOR_EMERALD)
    ]
    
    row_y = table_y + 45
    for eid, junc, cause, sev, res, col in events_mock:
        draw.line([(table_x, row_y), (table_x + table_w, row_y)], fill=COLOR_CARD_BORDER, width=1)
        draw.text((table_x + 20, row_y + 20), eid, font=load_font("semibold", 14), fill=COLOR_TEXT_PRIMARY)
        draw.text((table_x + 180, row_y + 20), junc, font=load_font("regular", 14), fill=COLOR_TEXT_PRIMARY)
        draw.text((table_x + 360, row_y + 20), cause, font=load_font("regular", 14), fill=COLOR_TEXT_PRIMARY)
        
        draw_card(draw, table_x + 500, row_y + 15, 110, 30, bg_color=(241, 245, 249), border_color=COLOR_CARD_BORDER, radius=6)
        draw_centered_text(draw, sev, table_x + 555, row_y + 30, load_font("bold", 11), col)
        
        draw.text((table_x + 650, row_y + 20), res, font=load_font("semibold", 13), fill=COLOR_TEXT_SECONDARY)
        row_y += 65
        
    draw_card(draw, 1000, 750, 800, 160, bg_color=COLOR_BG, border_color=COLOR_CARD_BORDER, radius=12)
    draw.text((1030, 770), "Incident Exploration Details:", font=load_font("bold", 16), fill=COLOR_TEXT_PRIMARY)
    draw.text((1030, 800), "• Full audit log captures exact timestamp, operator IDs, and latency.", font=load_font("regular", 14), fill=COLOR_TEXT_SECONDARY)
    draw.text((1030, 830), "• Search bar filters historical records by ID, junction name, or zone.", font=load_font("regular", 14), fill=COLOR_TEXT_SECONDARY)
    draw.text((1030, 860), "• CSV export button compiles active queries for administrative review.", font=load_font("regular", 14), fill=COLOR_TEXT_SECONDARY)

    img.save(os.path.join(SLIDES_DIR, "slide_06.png"))

def make_slide_7():
    img = Image.new("RGB", (WIDTH, HEIGHT), COLOR_BG)
    draw = ImageDraw.Draw(img)
    draw_header_footer(draw, 7, "Wireframes: Analytics & AI Triage Copilot")
    
    # Left Card: Analytics Dashboard
    draw_card(draw, 80, 210, 840, 730, radius=20)
    draw.text((120, 250), "Recharts Analytical Intelligence", font=load_font("bold", 24), fill=COLOR_TEXT_PRIMARY)
    
    cx, cy = 250, 450
    draw.ellipse([(cx - 90, cy - 90), (cx + 90, cy + 90)], outline=COLOR_CARD_BORDER, width=2)
    draw.pieslice([(cx - 90, cy - 90), (cx + 90, cy + 90)], start=0, end=120, fill=COLOR_RED)
    draw.pieslice([(cx - 90, cy - 90), (cx + 90, cy + 90)], start=120, end=240, fill=COLOR_AMBER)
    draw.pieslice([(cx - 90, cy - 90), (cx + 90, cy + 90)], start=240, end=310, fill=COLOR_BLUE)
    draw.pieslice([(cx - 90, cy - 90), (cx + 90, cy + 90)], start=310, end=360, fill=COLOR_EMERALD)
    draw.ellipse([(cx - 45, cy - 45), (cx + 45, cy + 45)], fill=COLOR_CARD_BG)
    draw_centered_text(draw, "128", cx, cy - 10, load_font("bold", 24), COLOR_TEXT_PRIMARY)
    draw_centered_text(draw, "Events", cx, cy + 22, load_font("regular", 12), COLOR_TEXT_SECONDARY)
    
    legend_y = 380
    for label, col in [("Critical", COLOR_RED), ("High", COLOR_AMBER), ("Medium", COLOR_BLUE), ("Low", COLOR_EMERALD)]:
        draw.ellipse([(390, legend_y), (405, legend_y + 15)], fill=col)
        draw.text((420, legend_y - 2), label, font=load_font("semibold", 14), fill=COLOR_TEXT_PRIMARY)
        legend_y += 35
        
    bar_x, bar_y, bar_w, bar_h = 120, 580, 760, 220
    draw_card(draw, bar_x, bar_y, bar_w, bar_h, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER, radius=12)
    
    draw.text((bar_x + 20, bar_y + 15), "INCIDENT CAUSE FREQUENCY", font=load_font("semibold", 12), fill=COLOR_TEXT_SECONDARY)
    
    causes = ["Accident", "Water Log", "Construction", "Breakdown", "Protest"]
    counts = [45, 32, 28, 16, 7]
    max_count = 50
    
    for i in range(len(causes)):
        bx = bar_x + 50 + i * 140
        bh = int((counts[i]/max_count) * 120)
        by = bar_y + 170 - bh
        draw.rectangle([(bx, by), (bx + 60, bar_y + 170)], fill=COLOR_BLUE)
        draw_centered_text(draw, str(counts[i]), bx + 30, by - 15, load_font("bold", 12), COLOR_TEXT_PRIMARY)
        draw_centered_text(draw, causes[i], bx + 30, bar_y + 185, load_font("regular", 12), COLOR_TEXT_SECONDARY)

    # Right Card: AI Triage Copilot
    draw_card(draw, 960, 210, 880, 730, radius=20)
    draw.text((1000, 250), "AI Copilot & Incident Simulation Form", font=load_font("bold", 24), fill=COLOR_TEXT_PRIMARY)
    
    form_x, form_y, form_w, form_h = 1000, 300, 800, 280
    draw_card(draw, form_x, form_y, form_w, form_h, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER, radius=12)
    
    fields = [
        ("EVENT TYPE", "unplanned"),
        ("CAUSE", "accident"),
        ("ROAD CLOSURE", "requires_closure"),
        ("VEHICLE TYPE", "heavy_vehicle"),
        ("CORRIDOR", "Hosur Road"),
        ("ZONE", "South Zone 1"),
        ("JUNCTION", "SilkBoardJunc"),
        ("COORDINATES", "12.9176, 77.6244")
    ]
    
    for i, (label, val) in enumerate(fields):
        row = i // 2
        col = i % 2
        fx = form_x + 30 + col * 380
        fy = form_y + 20 + row * 60
        draw.text((fx, fy), label, font=load_font("bold", 11), fill=COLOR_TEXT_MUTED)
        
        draw_card(draw, fx, fy + 18, 350, 35, bg_color=COLOR_BG, border_color=COLOR_CARD_BORDER, radius=6)
        draw.text((fx + 15, fy + 26), val, font=load_font("semibold", 13), fill=COLOR_TEXT_PRIMARY)
        
    btn_x, btn_y, btn_w, btn_h = form_x + 30, form_y + 225, 740, 42
    draw_card(draw, btn_x, btn_y, btn_w, btn_h, bg_color=COLOR_BLUE, border_color=COLOR_BLUE, radius=8)
    draw_centered_text(draw, "RUN AI INCIDENT TRIAGE", btn_x + btn_w/2, btn_y + btn_h/2, load_font("bold", 14), (255, 255, 255))
    
    verdict_x, verdict_y, verdict_w, verdict_h = 1000, 600, 800, 310
    draw_card(draw, verdict_x, verdict_y, verdict_w, verdict_h, bg_color=(239, 246, 255), border_color=(191, 219, 254), border_width=2, radius=16)
    
    draw.text((verdict_x + 30, verdict_y + 25), "AI INCIDENT VERDICT", font=load_font("bold", 16), fill=COLOR_BLUE)
    
    score_cx, score_cy = verdict_x + 130, verdict_y + 160
    draw.ellipse([(score_cx - 60, score_cy - 60), (score_cx + 60, score_cy + 60)], fill=COLOR_CARD_BG, outline=COLOR_RED, width=3)
    draw_centered_text(draw, "75%", score_cx, score_cy - 10, load_font("bold", 32), COLOR_RED)
    draw_centered_text(draw, "Criticality", score_cx, score_cy + 20, load_font("semibold", 12), COLOR_TEXT_SECONDARY)
    
    text_x = verdict_x + 230
    draw.text((text_x, verdict_y + 70), "IMPACT BAND:", font=load_font("semibold", 13), fill=COLOR_TEXT_SECONDARY)
    draw_card(draw, text_x + 115, verdict_y + 63, 100, 30, bg_color=COLOR_RED, border_color=COLOR_RED, radius=6)
    draw_centered_text(draw, "CRITICAL", text_x + 165, verdict_y + 78, load_font("bold", 11), (255, 255, 255))
    
    draw.text((text_x, verdict_y + 115), "ROAD CLOSURE:", font=load_font("semibold", 13), fill=COLOR_TEXT_SECONDARY)
    draw.text((text_x + 130, verdict_y + 115), "Highly Probable (84%)", font=load_font("bold", 14), fill=COLOR_TEXT_PRIMARY)
    
    draw.text((text_x, verdict_y + 155), "DISPATCH REC:", font=load_font("semibold", 13), fill=COLOR_TEXT_SECONDARY)
    draw.text((text_x + 130, verdict_y + 155), "15 Officers, 30 Barricades", font=load_font("bold", 14), fill=COLOR_TEXT_PRIMARY)
    
    draw.text((text_x, verdict_y + 195), "DIVERSION:", font=load_font("semibold", 13), fill=COLOR_TEXT_SECONDARY)
    draw_card(draw, text_x + 130, verdict_y + 188, 120, 30, bg_color=COLOR_EMERALD, border_color=COLOR_EMERALD, radius=6)
    draw_centered_text(draw, "MANDATORY", text_x + 190, verdict_y + 203, load_font("bold", 11), (255, 255, 255))
    
    draw.text((text_x, verdict_y + 245), "AI Insights: Unplanned heavy vehicle accident on Hosur Rd corridor.", font=load_font("bold", 13), fill=COLOR_TEXT_PRIMARY)
    draw.text((text_x, verdict_y + 268), "Junction SilkBoardJunc scores dynamic congestion index of 0.82.", font=load_font("regular", 13), fill=COLOR_TEXT_SECONDARY)

    img.save(os.path.join(SLIDES_DIR, "slide_07.png"))

def make_slide_8():
    img = Image.new("RGB", (WIDTH, HEIGHT), COLOR_BG)
    draw = ImageDraw.Draw(img)
    draw_header_footer(draw, 8, "System Architecture & Data Flow Diagram")
    
    draw_card(draw, 80, 210, 1760, 730, radius=24)
    draw.text((120, 260), "Decoupled 3-Tier Enterprise Infrastructure Layout", font=load_font("bold", 24), fill=COLOR_TEXT_PRIMARY)
    
    tiers = [
        ("1. CLIENT LAYER", "Next.js App Router (React 19)\nLeaflet Map HUD Engine\nRecharts Analytics Module\nLocal Storage State Controller", COLOR_BLUE),
        ("2. BACKEND API", "FastAPI Asynchronous Engine\nPydantic Input Validation\nRouter Endpoint Handlers\nDatabase Client Wrapper", COLOR_BLUE),
        ("3. MODEL / RULES", "CatBoost Classifier Engine\nOrdered Target Encoding\nContinuous Severity Scoring\nStatic Resource Rules Matrix", COLOR_AMBER),
        ("4. DATABASE TIER", "Supabase PostgreSQL (Cloud)\nSQLite In-Process (Fallback)\nJunction Running Score UPSERT\nSystem Audit Trail Loggers", COLOR_EMERALD)
    ]
    
    t_coords = [
        (130, 360, 340, 480),
        (560, 360, 340, 480),
        (990, 360, 340, 480),
        (1420, 360, 340, 480)
    ]
    
    for i, (title, contents, col) in enumerate(tiers):
        tx, ty, tw, th = t_coords[i]
        draw_card(draw, tx, ty, tw, th, bg_color=COLOR_CARD_BG, border_color=col, border_width=2, radius=16)
        
        draw.rectangle([(tx+2, ty+2), (tx+tw-2, ty+60)], fill=(241, 245, 249))
        draw_centered_text(draw, title, tx + tw/2, ty + 31, load_font("bold", 16), col)
        
        y = ty + 90
        for line in contents.split('\n'):
            draw.ellipse([(tx + 30, y + 8), (tx + 38, y + 16)], fill=col)
            draw_wrapped_text(draw, line, tx + 55, y, tw - 80, load_font("semibold", 14), COLOR_TEXT_PRIMARY)
            y += 85
            
    draw.line([(470, 600), (560, 600)], fill=COLOR_TEXT_MUTED, width=4)
    draw.polygon([(560, 600), (545, 590), (545, 610)], fill=COLOR_TEXT_MUTED)
    draw_centered_text(draw, "REST HTTP\nJSON Payload", 515, 560, load_font("regular", 11), COLOR_TEXT_SECONDARY)
    
    draw.line([(900, 600), (990, 600)], fill=COLOR_TEXT_MUTED, width=4)
    draw.polygon([(990, 600), (975, 590), (975, 610)], fill=COLOR_TEXT_MUTED)
    draw_centered_text(draw, "In-Process\nInference", 945, 560, load_font("regular", 11), COLOR_TEXT_SECONDARY)
    
    draw.line([(730, 840), (730, 900), (1590, 900), (1590, 840)], fill=COLOR_TEXT_MUTED, width=4)
    draw.polygon([(1590, 840), (1580, 855), (1600, 855)], fill=COLOR_TEXT_MUTED)
    draw_centered_text(draw, "Double-Tier DB sync (Supabase Primary / SQLite local Fallback)", 1160, 920, load_font("bold", 13), COLOR_TEXT_SECONDARY)

    img.save(os.path.join(SLIDES_DIR, "slide_08.png"))

def make_slide_9():
    img = Image.new("RGB", (WIDTH, HEIGHT), COLOR_BG)
    draw = ImageDraw.Draw(img)
    draw_header_footer(draw, 9, "Technology Stack & AI/ML Components")
    
    draw_card(draw, 80, 210, 1760, 730, radius=24)
    
    draw.text((120, 260), "Technological Infrastructure Stack", font=load_font("bold", 24), fill=COLOR_TEXT_PRIMARY)
    
    tech_stack = [
        ("Frontend Layer", "React 19  •  Next.js (App Router)  •  Tailwind CSS  •  Leaflet Maps  •  Recharts", COLOR_BLUE),
        ("Backend Layer", "FastAPI (Asynchronous REST)  •  Python 3.11  •  Uvicorn Server  •  Pydantic V2 Schema", COLOR_BLUE),
        ("Storage Layer", "Supabase PostgreSQL (AWS Cloud Database)  •  SQLite3 Local SQL Fallback Engine", COLOR_EMERALD),
        ("Deployment Layer", "Render Web Services (FastAPI Backend)  •  Vercel Edge Platform (Next.js client)", COLOR_EMERALD)
    ]
    
    cy = 310
    for label, desc, col in tech_stack:
        draw_card(draw, 120, cy, 780, 110, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER, radius=12)
        draw.rectangle([(120, cy + 10), (128, cy + 100)], fill=col)
        draw.text((150, cy + 20), label, font=load_font("bold", 18), fill=col)
        draw.text((150, cy + 55), desc, font=load_font("semibold", 16), fill=COLOR_TEXT_PRIMARY)
        cy += 140
        
    draw.text((960, 260), "CatBoost Classifier & Severity Calculations", font=load_font("bold", 24), fill=COLOR_TEXT_PRIMARY)
    
    draw_card(draw, 960, 310, 800, 530, bg_color=COLOR_CARD_BG, border_color=COLOR_AMBER, border_width=2, radius=16)
    
    draw_brain_icon(draw, 1020, 370)
    draw.text((1070, 345), "Why CatBoost Algorithm?", font=load_font("bold", 22), fill=COLOR_AMBER)
    
    cb_desc = "Our spatial incident dataset is heavily dominated by high-cardinality categorical variables (e.g. junction names, corridors, zones). CatBoost handles these natively using Ordered Target Encoding without the dimensionality explosion of one-hot methods. Symmetric decision trees yield predictions in under 10ms with robust resistance to overfitting."
    draw_wrapped_text(draw, cb_desc, 1070, 390, 650, load_font("regular", 14), COLOR_TEXT_SECONDARY)
    
    draw_card(draw, 1000, 540, 720, 120, bg_color=(254, 243, 199), border_color=(251, 191, 36), radius=12)
    draw.text((1030, 560), "CONTINUOUS SEVERITY SCORE (IMPACT INDEX) FORMULA:", font=load_font("bold", 13), fill=(180, 83, 9))
    
    formula_text = "Impact Score = 0.10*P(Low) + 0.40*P(Medium) + 0.75*P(High) + 1.00*P(Critical)"
    draw.text((1030, 600), formula_text, font=load_font("bold", 20), fill=COLOR_TEXT_PRIMARY)
    
    draw.text((1000, 690), "Validated Performance Matrix (backend/test_model.py):", font=load_font("bold", 15), fill=COLOR_TEXT_PRIMARY)
    
    perf_bullets = [
        "Inference Latency: 9.2 milliseconds (in-process model evaluation).",
        "Categorical handling: Preprocessor maps string coordinates to Bengaluru spatial zones.",
        "Resource counts are matched via local rule engine lookup in resource_map.json."
    ]
    y = 730
    for bullet in perf_bullets:
        draw.ellipse([(1000, y + 6), (1006, y + 12)], fill=COLOR_TEXT_SECONDARY)
        draw.text((1025, y - 2), bullet, font=load_font("regular", 14), fill=COLOR_TEXT_SECONDARY)
        y += 30

    img.save(os.path.join(SLIDES_DIR, "slide_09.png"))

def make_slide_10():
    img = Image.new("RGB", (WIDTH, HEIGHT), COLOR_BG)
    draw = ImageDraw.Draw(img)
    draw_header_footer(draw, 10, "Expected Impact, Scalability & Future Scope")
    
    draw_card(draw, 80, 210, 1760, 730, radius=24)
    
    cols = [
        ("1. EXPECTED IMPACT", [
            ("Rapid Dispatch Turnaround", "AI auto-allocation reduces incident assessment cycles from hours to seconds."),
            ("Emergency Priority Routing", "Predictive traffic hotspots flag bottlenecks ahead of time, ensuring clear lanes for emergency services."),
            ("Data-Driven Operations", "Junction severity indexes are maintained in dynamic database registers, letting administrators track seasonal bottlenecks."),
            ("Clearance Auditing", "Timestamp logging on SQLite/Supabase gives audit trails of when resources are deployed vs. congestion clearance.")
        ], COLOR_EMERALD),
        
        ("2. SYSTEM SCALABILITY", [
            ("Cloud-Native API Structure", "FastAPI backend containerizes easily via Docker, deploying across load-balanced AWS ECS nodes."),
            ("Multi-City Extensibility", "The database coordinates schema supports variable zones, corridors, and junctions, making it adaptable to other metros."),
            ("Decoupled Frontend HUD", "Next.js client runs on edge networks, rendering maps and graphs using client-side computations."),
            ("Offline Resilience", "Local fallback databases prevent control room blackouts even during cloud infrastructure server outages.")
        ], COLOR_BLUE),
        
        ("3. FUTURE ROADMAP", [
            ("IoT & CCTV Ingestion", "Integrate live video traffic streams to auto-detect vehicle breakdowns and accidents via Computer Vision."),
            ("Real-Time GPS Telemetry", "Incorporate live speed logs from commercial vehicle fleets to calculate dynamically weighted routing graphs."),
            ("Continuous Loop Sensors", "Deploy hardware counters at junctions to feed real-time speed and traffic volume variables into the model."),
            ("Voice Operator Copilot", "Incorporate Generative LLMs to let operators dispatch personnel using verbal commands.")
        ], COLOR_AMBER)
    ]
    
    col_w = 510
    gap = 60
    
    for i, (title, items, col) in enumerate(cols):
        cx = 140 + i * (col_w + gap)
        cy = 270
        
        draw.text((cx, cy), title, font=load_font("bold", 22), fill=col)
        draw.line([(cx, cy + 38), (cx + col_w, cy + 38)], fill=col, width=3)
        
        y = cy + 60
        for item_title, item_desc in items:
            draw_card(draw, cx, y, col_w, 110, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER, radius=10)
            draw.rectangle([(cx, y + 10), (cx + 6, y + 100)], fill=col)
            
            draw.text((cx + 20, y + 12), item_title, font=load_font("bold", 15), fill=COLOR_TEXT_PRIMARY)
            draw_wrapped_text(draw, item_desc, cx + 20, y + 36, col_w - 40, load_font("regular", 12), COLOR_TEXT_SECONDARY)
            y += 130

    img.save(os.path.join(SLIDES_DIR, "slide_10.png"))

def make_slide_11():
    img = Image.new("RGB", (WIDTH, HEIGHT), COLOR_BG)
    draw = ImageDraw.Draw(img)
    
    draw.text((80, 45), "Flipkart", font=load_font("bold", 30), fill=COLOR_BLUE)
    draw.text((205, 42), "|", font=load_font("regular", 32), fill=COLOR_TEXT_MUTED)
    draw.text((230, 46), "hackerearth", font=load_font("bold", 28), fill=(0, 0, 0))
    
    draw.text((1350, 46), "Flipkart", font=load_font("regular", 22), fill=COLOR_TEXT_PRIMARY)
    draw.text((1435, 46), "GRIDLOCK HACKATHON", font=load_font("bold", 22), fill=COLOR_TEXT_PRIMARY)
    badge_x, badge_y, badge_w, badge_h = 1690, 44, 150, 36
    draw.rounded_rectangle([(badge_x, badge_y), (badge_x + badge_w, badge_y + badge_h)], radius=8, fill=(254, 240, 138))
    draw_centered_text(draw, "solve. traffic", badge_x + badge_w/2, badge_y + badge_h/2, load_font("bold", 18), (113, 63, 18))
    
    draw_centered_text(draw, "THANK YOU", WIDTH/2, HEIGHT/2 - 120, load_font("bold", 96), COLOR_TEXT_PRIMARY)
    draw_centered_text(draw, "Smarter Traffic Operations. Safer Roads. Better Cities.", WIDTH/2, HEIGHT/2, load_font("semibold", 28), COLOR_BLUE)
    
    card_w, card_h = 600, 220
    cx, cy = WIDTH/2 - card_w/2, HEIGHT/2 + 100
    draw_card(draw, cx, cy, card_w, card_h, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER, radius=20)
    
    draw_centered_text(draw, "Presented By: Team Arohan", cx + card_w/2, cy + 40, load_font("bold", 20), COLOR_TEXT_PRIMARY)
    
    leader_text = "Team Leader: Nikita Pawar"
    members_text = "Jay Bhandarkar  •  Prajwal Kate  •  Riya Lad"
    draw_centered_text(draw, leader_text, cx + card_w/2, cy + 100, load_font("semibold", 18), COLOR_TEXT_SECONDARY)
    draw_centered_text(draw, members_text, cx + card_w/2, cy + 145, load_font("semibold", 18), COLOR_TEXT_SECONDARY)
    
    draw.rectangle([(0, 1072), (640, 1080)], fill=COLOR_BLUE)
    draw.rectangle([(640, 1072), (1280, 1080)], fill=COLOR_EMERALD)
    draw.rectangle([(1280, 1072), (1920, 1080)], fill=COLOR_AMBER)

    img.save(os.path.join(SLIDES_DIR, "slide_11.png"))

# ----------------- MAIN COMPILATION AND EXPORT -----------------

def generate_all_images():
    print("[1/2] Generating slide PNG images...")
    make_slide_1()
    make_slide_2()
    make_slide_3()
    make_slide_4()
    make_slide_5()
    make_slide_6()
    make_slide_7()
    make_slide_8()
    make_slide_9()
    make_slide_10()
    make_slide_11()
    print("All slide images generated in backend/generated_slides/")

def compile_to_pptx():
    print("[2/2] Compiling slide PNGs into PowerPoint PPTX...")
    prs = Presentation()
    
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    for i in range(1, 12):
        slide = prs.slides.add_slide(blank_slide_layout)
        slide_img_path = os.path.join(SLIDES_DIR, f"slide_{i:02d}.png")
        
        slide.shapes.add_picture(slide_img_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        
    prs.save(PPTX_PATH)
    print(f"Presentation PPTX saved successfully at: {PPTX_PATH}")

if __name__ == "__main__":
    generate_all_images()
    compile_to_pptx()
